"""
Build AFFG Managed Device Strategy presentation.
Technijian Brand Guide 2026. Slide size: 13.333" x 7.5" (PowerPoint widescreen 16:9).
All positions calibrated for 7.5" height (brand spec was 5.63" - scaled up).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
CORE_ORANGE   = RGBColor(0xF6, 0x7D, 0x4B)
CORE_BLUE     = RGBColor(0x00, 0x6D, 0xB6)
TEAL          = RGBColor(0x1E, 0xAA, 0xC8)
DARK_CHARCOAL = RGBColor(0x1A, 0x1A, 0x2E)
BRAND_GREY    = RGBColor(0x59, 0x59, 0x5B)
OFF_WHITE     = RGBColor(0xF8, 0xF9, 0xFA)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY    = RGBColor(0xE9, 0xEC, 0xEF)
ROW_ALT       = RGBColor(0xF1, 0xF3, 0xF5)

# ── Logo Paths ──
# White logo hi-res transparent - for blue header bars
LOGO_WHITE = r'C:\vscode\tech-branding\tech-branding\assets\logos\png\technijian-logo-reverse-white-5000x1667-transparent.png'
# Full-color transparent - for light/white backgrounds
LOGO_COLOR = r'C:\vscode\tech-branding\tech-branding\assets\logos\png\technijian-logo-full-color-2400x502.png'
# Dark-bg transparent (colored dots + white text) - for dark charcoal slides
LOGO_DARK  = r'C:\vscode\tech-branding\tech-branding\assets\logos\png\technijian-logo-dark-bg-transparent.png'

# ── Slide Dimensions ──
SW = 12192000   # 13.333"
SH = 6858000    # 7.5"

# ── Layout Constants (calibrated for 7.5" slide height) ──
HEADER_H     = Inches(1.1)       # blue header bar
# Logo sizes calculated from actual image aspect ratios
# reverse-white-5000x1667 = 3:1 ratio (includes tagline)
LOGO_WHITE_W = Inches(2.4)       # white logo in header
LOGO_WHITE_H = Inches(0.8)       # 2.4/3.0 = 0.8"
# dark-bg-transparent 900x160 = 5.62:1 ratio
LOGO_DARK_W  = Inches(2.6)       # dark-bg logo on title slide
LOGO_DARK_H  = Inches(0.46)      # 2.6/5.62 = 0.46"
TITLE_Y      = Inches(1.25)      # slide title Y position
CONTENT_Y    = Inches(2.1)       # where cards/tables start
FOOTER_Y     = Inches(6.9)       # footer line
FOOTER_TXT_Y = Inches(6.95)      # footer text
CONTENT_H    = FOOTER_Y - CONTENT_Y  # usable content height ~4.8"

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH

# ── Helpers ──

def new_slide():
    """Create a blank slide with all default placeholders removed."""
    s = prs.slides.add_slide(prs.slide_layouts[0])
    # Remove default Title and Subtitle placeholder shapes that show as boxes behind content
    for ph in list(s.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    return s

def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    return s

def txt(slide, l, t, w, h, text, sz, bold=False, color=DARK_CHARCOAL, align=PP_ALIGN.LEFT):
    from pptx.oxml.ns import qn
    s = slide.shapes.add_textbox(l, t, w, h)
    # Force no fill - remove any default spAutoFit and set noFill explicitly
    sp = s._element
    sp_pr = sp.find(qn('a:spPr')) if sp.find(qn('a:spPr')) is not None else sp.makeelement(qn('a:spPr'), {})
    no_fill = sp_pr.makeelement(qn('a:noFill'), {})
    # Remove existing fill
    for child in list(sp_pr):
        if child.tag in [qn('a:solidFill'), qn('a:noFill'), qn('a:gradFill'), qn('a:pattFill')]:
            sp_pr.remove(child)
    sp_pr.append(no_fill)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = sz; r.font.bold = bold
    r.font.color.rgb = color; r.font.name = 'Open Sans'
    return s

def mtxt(slide, l, t, w, h, lines, sz=Pt(13), color=BRAND_GREY, sp=Pt(5)):
    from pptx.oxml.ns import qn
    s = slide.shapes.add_textbox(l, t, w, h)
    sp_el = s._element
    sp_pr = sp_el.find(qn('a:spPr')) if sp_el.find(qn('a:spPr')) is not None else sp_el.makeelement(qn('a:spPr'), {})
    no_fill = sp_pr.makeelement(qn('a:noFill'), {})
    for child in list(sp_pr):
        if child.tag in [qn('a:solidFill'), qn('a:noFill'), qn('a:gradFill'), qn('a:pattFill')]:
            sp_pr.remove(child)
    sp_pr.append(no_fill)
    tf = s.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = sp; p.space_after = sp
        r = p.add_run()
        r.text = line; r.font.size = sz; r.font.color.rgb = color; r.font.name = 'Open Sans'
    return s

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def chrome(slide, title_text, pg):
    """Blue header bar with logo, title, footer."""
    rect(slide, 0, 0, SW, HEADER_H, CORE_BLUE)
    slide.shapes.add_picture(LOGO_WHITE, Inches(0.4), Inches(0.15), LOGO_WHITE_W, LOGO_WHITE_H)
    txt(slide, Inches(11.5), Inches(0.35), Inches(1.2), Inches(0.4), str(pg), Pt(11), color=WHITE, align=PP_ALIGN.RIGHT)
    txt(slide, Inches(0.6), TITLE_Y, Inches(12), Inches(0.7), title_text, Pt(28), True, CORE_BLUE)
    rect(slide, Inches(0.6), FOOTER_Y, Inches(12), Inches(0.01), LIGHT_GREY)
    txt(slide, Inches(0.6), FOOTER_TXT_Y, Inches(4), Inches(0.3), 'technijian.com', Pt(9), color=BRAND_GREY)
    txt(slide, Inches(8), FOOTER_TXT_Y, Inches(4.5), Inches(0.3),
        'AFFG Managed Device Strategy  |  April 2026', Pt(9), color=BRAND_GREY, align=PP_ALIGN.RIGHT)

def card3(slide, idx, title_text, body_lines, citation):
    card_w = Inches(3.8); gap = Inches(0.2)
    starts = [Inches(0.6), Inches(0.6) + card_w + gap, Inches(0.6) + 2*(card_w + gap)]
    x = starts[idx]
    card_h = FOOTER_Y - CONTENT_Y - Inches(0.15)
    rect(slide, x, CONTENT_Y, card_w, card_h, OFF_WHITE)
    txt(slide, x+Inches(0.3), CONTENT_Y+Inches(0.2), card_w-Inches(0.6), Inches(0.55), title_text, Pt(16), True, CORE_BLUE)
    rect(slide, x+Inches(0.3), CONTENT_Y+Inches(0.8), Inches(1.0), Inches(0.04), CORE_ORANGE)
    mtxt(slide, x+Inches(0.3), CONTENT_Y+Inches(0.95), card_w-Inches(0.6), card_h-Inches(1.5), body_lines, Pt(12), BRAND_GREY, Pt(4))
    txt(slide, x+Inches(0.3), CONTENT_Y+card_h-Inches(0.45), card_w-Inches(0.6), Inches(0.4), citation, Pt(9), True, TEAL)

def card4(slide, idx, title_text, body_lines, citation):
    card_w = Inches(2.85); gap = Inches(0.15)
    x = Inches(0.6) + idx * (card_w + gap)
    card_h = FOOTER_Y - CONTENT_Y - Inches(0.15)
    rect(slide, x, CONTENT_Y, card_w, card_h, OFF_WHITE)
    txt(slide, x+Inches(0.2), CONTENT_Y+Inches(0.15), card_w-Inches(0.4), Inches(0.5), title_text, Pt(14), True, CORE_BLUE)
    rect(slide, x+Inches(0.2), CONTENT_Y+Inches(0.65), Inches(0.8), Inches(0.04), CORE_ORANGE)
    mtxt(slide, x+Inches(0.2), CONTENT_Y+Inches(0.8), card_w-Inches(0.4), card_h-Inches(1.3), body_lines, Pt(10.5), BRAND_GREY, Pt(3))
    txt(slide, x+Inches(0.2), CONTENT_Y+card_h-Inches(0.4), card_w-Inches(0.4), Inches(0.35), citation, Pt(8), True, TEAL)

def table_slide(slide, title_text, pg, headers, rows):
    chrome(slide, title_text, pg)
    row_l = Inches(0.6); row_w = Inches(12)
    hdr_top = CONTENT_Y; hdr_h = Inches(0.4)
    rect(slide, row_l, hdr_top, row_w, hdr_h, CORE_BLUE)
    nc = len(headers)
    if nc == 3:
        offs = [Inches(0.75), Inches(3.3), Inches(8.0)]
        ws   = [Inches(2.4), Inches(4.5), Inches(4.2)]
    elif nc == 4:
        offs = [Inches(0.75), Inches(2.7), Inches(5.7), Inches(9.3)]
        ws   = [Inches(1.8), Inches(2.8), Inches(3.4), Inches(2.8)]
    else:
        offs = [Inches(0.75), Inches(3.3), Inches(8.0)]
        ws   = [Inches(2.4), Inches(4.5), Inches(4.2)]
    for j, h in enumerate(headers):
        txt(slide, offs[j], hdr_top+Inches(0.07), ws[j], Inches(0.28), h, Pt(11), True, WHITE)
    nr = len(rows)
    avail = FOOTER_Y - (hdr_top + hdr_h + Inches(0.05))
    rh = int(avail / nr); rh = min(rh, Inches(0.6))
    fsz = Pt(10.5) if rh >= Inches(0.5) else Pt(9.5)
    cur = hdr_top + hdr_h + Inches(0.02)
    for i, rd in enumerate(rows):
        bg = WHITE if i % 2 == 0 else ROW_ALT
        rect(slide, row_l, cur, row_w, rh, bg)
        for j, (ct, b, c) in enumerate(rd):
            txt(slide, offs[j], cur+Inches(0.05), ws[j], rh-Inches(0.1), ct, fsz, b, c)
        cur += rh

def bullet_slide(slide, title_text, pg, bullets):
    chrome(slide, title_text, pg)
    n = len(bullets)
    avail = FOOTER_Y - CONTENT_Y - Inches(0.1)
    sp = int(avail / n); sp = min(sp, Inches(0.55))
    for i, (dc, bt) in enumerate(bullets):
        y = CONTENT_Y + i * sp
        rect(slide, Inches(0.8), y+Inches(0.08), Inches(0.1), Inches(0.1), dc)
        txt(slide, Inches(1.1), y, Inches(11.3), Inches(0.45), bt, Pt(13), color=BRAND_GREY)


# ═════════════════════════════════════════════
# SLIDE 1: TITLE (Dark Charcoal)
# ═════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, SW, SH, DARK_CHARCOAL)
s.shapes.add_picture(LOGO_WHITE, Inches(0.6), Inches(0.4), Inches(2.6), Inches(0.87))
txt(s, Inches(0.6), Inches(2.0), Inches(8.8), Inches(1.0), 'AFFG Managed Device Strategy', Pt(40), True, WHITE)
txt(s, Inches(0.6), Inches(3.2), Inches(8.8), Inches(0.7), 'Company Laptops + CloudBrink ZTNA', Pt(24), True, WHITE)
txt(s, Inches(0.6), Inches(4.2), Inches(8.8), Inches(0.8),
    'Same FINRA and SEC compliance posture without VDI \u2014 enforced by managed endpoints, CloudBrink zero-trust access, and Intune device management',
    Pt(15), color=LIGHT_GREY)
rect(s, 0, Inches(6.85), SW, Inches(0.08), CORE_ORANGE)
txt(s, Inches(7.0), Inches(6.3), Inches(5.5), Inches(0.4), 'technology as a solution', Pt(11), color=TEAL, align=PP_ALIGN.RIGHT)
txt(s, Inches(0.6), Inches(7.05), Inches(4), Inches(0.3), 'Prepared for AFFG  |  April 16, 2026', Pt(10), color=BRAND_GREY)
notes(s, """TALKING POINTS:
- Proposes replacing VDI with company-managed laptops and CloudBrink ZTNA for all 9 AFFG users.
- Same FINRA/SEC compliance posture. Every control maps to a specific citation.
- Users get native laptop performance, work securely from any location.
- All 9 users receive company laptop + company phone, both Intune-managed.
- No additional M365 licensing - Intune included in E3.""")


# ═════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ═════════════════════════════════════════════
s = new_slide()
chrome(s, 'Executive Summary', 2)
card3(s, 0, 'The Shift: VDI To Managed Devices',
      ['Instead of hosting desktops in a datacenter, AFFG provides company-owned laptops and phones to all 9 users.',
       'Data protection moves from "data never leaves the datacenter" to "the device is the controlled boundary."',
       'Same regulatory outcome \u2014 different enforcement model, lower complexity, better user experience.'],
      'Reg S-P \u00a7248.30(a); FINRA 3110(a)')
card3(s, 1, 'The Challenge: Access From Anywhere',
      ['VDI had a natural boundary: data stayed in the datacenter. Managed laptops travel.',
       'IP whitelisting alone fails when users work from hotels, client sites, or home.',
       'CloudBrink ZTNA replaces the fixed-IP boundary with a device-identity + posture boundary that works from any network.'],
      'Reg S-P \u00a7248.30(a)(3); NIST AC-17')
card3(s, 2, 'The Design: Layered Enforcement',
      ['The control stack layers Technijian services:',
       '\u2022 Company laptops (Intune-managed)',
       '\u2022 Company phones (Intune MDM/MAM)',
       '\u2022 CloudBrink ZTNA (zero-trust access)',
       '\u2022 MyAudit endpoint DLP',
       '\u2022 SSO/2FA gateway',
       '\u2022 Veeam tenant backup + monthly assessment'],
      'Full stack: Reg S-P + FINRA 3110 + 17a-4')
notes(s, """TALKING POINTS:
- THE SHIFT: Datacenter VDI to company-owned managed devices. Endpoint becomes the controlled boundary.
- THE CHALLENGE: How to protect data when laptop travels. CloudBrink verifies device + user regardless of network.
- THE DESIGN: Seven enforcement layers. Each maps to a regulatory citation. No VDI hosting cost. Intune in E3.""")


# ═════════════════════════════════════════════
# SLIDE 3: CONTROL-TO-CITATION MAP
# ═════════════════════════════════════════════
s = new_slide()
table_slide(s, 'Control-to-Citation Map', 3,
    ['Design Component', 'Control Objective', 'Citation(s)'],
    [
        [('Intune-managed company laptops', True, DARK_CHARCOAL),
         ('All 9 users on company-owned, encrypted, policy-enforced devices.', False, BRAND_GREY),
         ('Reg S-P \u00a7248.30(a)(1)-(3); FINRA 3110(a)', False, BRAND_GREY)],
        [('Conditional Access + legacy auth block', True, DARK_CHARCOAL),
         ('Entra ID restricts M365 to compliant devices. IMAP/POP3/SMTP blocked.', False, BRAND_GREY),
         ('Reg S-P \u00a7248.30(a)(3); NIST AC-3, IA-2(1)', False, BRAND_GREY)],
        [('CloudBrink ZTNA', True, DARK_CHARCOAL),
         ('Zero-trust access replaces IP whitelist. Device posture verified per session.', False, BRAND_GREY),
         ('Reg S-P \u00a7248.30(a)(3); NIST AC-17, SC-7', False, BRAND_GREY)],
        [('Technijian MyAudit (endpoint DLP)', True, DARK_CHARCOAL),
         ('Block USB, local download, print, clipboard, screen capture.', False, BRAND_GREY),
         ('Reg S-P \u00a7248.30(a)(3); NIST AC-19, MP-7', False, BRAND_GREY)],
        [('Technijian SSO / 2FA gateway', True, DARK_CHARCOAL),
         ('Enforced, logged MFA on custodian portals and third-party SaaS.', False, BRAND_GREY),
         ('Reg S-P \u00a7248.30 (2024); NIST IA-2(1)', False, BRAND_GREY)],
        [('Intune MDM/MAM (company phones)', True, DARK_CHARCOAL),
         ('Phone enrollment, app containerization, remote wipe.', False, BRAND_GREY),
         ('Reg S-P \u00a7248.30(a)(1)-(3); NIST AC-19, MP-6', False, BRAND_GREY)],
        [('Technijian Veeam Backup for M365', True, DARK_CHARCOAL),
         ('Immutable, non-rewriteable backup of full M365 tenant.', False, BRAND_GREY),
         ('SEC Rule 17a-4(b),(f); FINRA 4511', False, BRAND_GREY)],
        [('Technijian monthly assessment', True, DARK_CHARCOAL),
         ('Detect drift; attested evidence of ongoing control operation.', False, BRAND_GREY),
         ('FINRA 3110(c); FINRA 3120', False, BRAND_GREY)],
    ])
notes(s, """TALKING POINTS:
- Master reference. Every component maps to a FINRA/SEC citation.
- If it can't be tied to a citation, it's not in the proposal.
- Conditional Access + legacy auth block is HOW we restrict M365 to company devices.
- Veeam unchanged from VDI proposal - backs up tenant, not endpoint.""")


# ═════════════════════════════════════════════
# SLIDE 4: MANAGED ENDPOINTS
# ═════════════════════════════════════════════
s = new_slide()
chrome(s, 'Managed Endpoints', 4)
card3(s, 0, 'Intune Enrollment',
      ['All 9 company laptops auto-enroll via Autopilot.',
       'Policies: BitLocker encryption, Windows Defender ATP, firewall, OS patching, screen lock.',
       'Non-compliant device loses M365 access within minutes.',
       'Intune included in M365 E3 \u2014 no additional cost.'],
      'Reg S-P \u00a7248.30(a)(1)-(3); NIST CM-8')
card3(s, 1, 'Conditional Access',
      ['M365 access requires all four factors:',
       '1. Correct user identity',
       '2. MFA authentication',
       '3. Intune-compliant device',
       '4. CloudBrink ZTNA session active',
       'Personal laptop with valid credentials + MFA still blocked.'],
      'Reg S-P \u00a7248.30(a); NIST AC-3, IA-2(1)')
card3(s, 2, 'Examiner Optic',
      ['Intune portal: every device with model, OS, encryption, compliance state, last check-in.',
       'Auditable, exportable device inventory.',
       'Lost laptop: remote wipe, BitLocker key escrowed.',
       'Terminated employee: auto-lock, wipe in minutes.'],
      'Reg S-P \u00a7248.30(a); FINRA 3110(a)')
notes(s, """TALKING POINTS:
- Enrollment: Autopilot, BitLocker enforced before data access.
- Conditional Access: four-factor boundary. Stolen credentials + MFA still blocked without compliant device.
- Examiner Optic: Intune portal is cleaner evidence than VDI session logs.
- Lost laptop: BitLocker protects data at rest + remote wipe. Equivalent to VDI "no data on device."
- No additional licensing: Intune in existing E3.""")


# ═════════════════════════════════════════════
# SLIDE 5: CONDITIONAL ACCESS CONFIGURATION
# ═════════════════════════════════════════════
s = new_slide()
chrome(s, 'Conditional Access Configuration', 5)
card3(s, 0, 'Policy 1: Require Compliant Device',
      ['Entra ID \u2192 Conditional Access \u2192 Policies',
       'Target: Exchange Online, SharePoint/OneDrive, Teams',
       'Applies to: All 9 users, all platforms',
       'Grant: "Require device to be marked as compliant"',
       'Result: unmanaged device \u2192 access denied. Blocks personal laptops, phones, unmanaged browsers.'],
      'Entra ID \u2192 Conditional Access \u2192 Grant')
card3(s, 1, 'Policy 2: Block Legacy Auth',
      ['CRITICAL: IMAP, POP3, SMTP basic auth bypass Conditional Access entirely.',
       'Without this, personal mail client with basic auth pulls all email.',
       'Target: "Other clients" + "Exchange ActiveSync" \u2192 Block.',
       'Also disable per-user legacy auth in Entra ID.',
       'Most commonly missed step in M365 lockdown.'],
      'Entra ID \u2192 Conditional Access \u2192 Block')
card3(s, 2, 'Policy 3: App Protection (Mobile)',
      ['Intune \u2192 App Protection Policies',
       'Target: Outlook, OneDrive, Teams, Edge',
       'Block: copy/paste to personal apps, save-as to personal storage, screenshots.',
       'Require: PIN/biometric for corporate apps.',
       'Auto-wipe: on jailbreak or X days offline.'],
      'Intune \u2192 App Protection Policies')
notes(s, """TALKING POINTS - CRITICAL CONFIGURATION:
POLICY 1: Entra ID CA policy targeting M365 apps. Grant = "Require compliant device." Blocks OWA from personal laptop, Outlook from personal phone, OneDrive web from unmanaged browser.
POLICY 2 (CRITICAL): Legacy auth bypasses CA entirely. Without this, IMAP/POP3 on personal phone pulls all email. Must block "Other clients" + "Exchange ActiveSync."
POLICY 3: Even on compliant phones, data stays containerized. No copy to WhatsApp, no save to personal storage.
IMPLEMENTATION ORDER: 1) Set up compliance policies, 2) Enroll devices, 3) Enable legacy auth block, 4) Enable CA in report-only 48hrs, 5) Enforce, 6) Deploy MAM, 7) Test.""")


# ═════════════════════════════════════════════
# SLIDE 6: CLOUDBRINK ZTNA
# ═════════════════════════════════════════════
s = new_slide()
chrome(s, 'CloudBrink ZTNA', 6)
card3(s, 0, 'How It Works',
      ['CloudBrink agent on every managed laptop.',
       'Verifies device is Intune-enrolled, compliant, user authenticated before any connection.',
       'End-to-end encrypted traffic through CloudBrink edge.',
       'Per-application access \u2014 no lateral movement.',
       'Works from office, hotel, client site, or home.'],
      'Reg S-P \u00a7248.30(a)(3); NIST AC-17, SC-7')
card3(s, 1, 'Why Not Just IP Whitelist',
      ['IP whitelisting works when all 9 users are in the office.',
       'Travel = locked out, or exceptions that weaken the boundary.',
       'CloudBrink: is this a known, managed, compliant device with a verified user?',
       'Stronger than IP whitelist \u2014 validates device, not network.'],
      'Reg S-P \u00a7248.30(a)(3); NIST AC-3, AC-17')
card3(s, 2, 'What It Replaces',
      ['\u2022 Replaces: office-IP whitelist for M365 and portals',
       '\u2022 Replaces: traditional VPN (split-tunnel risks)',
       '\u2022 Enforces: device posture check per connection',
       '\u2022 Enforces: per-app micro-segmented access',
       '\u2022 Logs: every session (user, device, app, duration)',
       '\u2022 Works: from any network, no config changes'],
      'Audit: SEC 17a-4(b)(4); FINRA 4511')
notes(s, """TALKING POINTS:
- CloudBrink is the key change enabling managed devices for a traveling workforce.
- HOW IT WORKS: Agent verifies device + user before any connection. Encrypted, per-app access.
- WHY NOT IP WHITELIST: Works in office, breaks when traveling. CloudBrink eliminates this.
- WHAT IT REPLACES: IP whitelist, VPN. Security is stronger - validates device, not network.""")


# ═════════════════════════════════════════════
# SLIDE 7: ACCESS BOUNDARY COMPARISON
# ═════════════════════════════════════════════
s = new_slide()
table_slide(s, 'VDI + IP Whitelist vs. Managed Device + CloudBrink', 7,
    ['Factor', 'VDI + IP Whitelist', 'Managed + CloudBrink', 'Impact'],
    [
        [('Data at rest', True, DARK_CHARCOAL), ('Data in datacenter VM', False, BRAND_GREY),
         ('BitLocker-encrypted laptop', False, BRAND_GREY), ('Equivalent', False, TEAL)],
        [('Access control', True, DARK_CHARCOAL), ('Office IP + MFA + endpoint', False, BRAND_GREY),
         ('ZTNA + MFA + Intune compliance', False, BRAND_GREY), ('Stronger', False, TEAL)],
        [('Travel / remote', True, DARK_CHARCOAL), ('Locked out or VPN exception', False, BRAND_GREY),
         ('Works from any network', False, BRAND_GREY), ('Eliminates risk', False, TEAL)],
        [('Data exfiltration', True, DARK_CHARCOAL), ('VDI session + MyAudit', False, BRAND_GREY),
         ('MyAudit on managed laptop', False, BRAND_GREY), ('Equivalent', False, TEAL)],
        [('Device loss', True, DARK_CHARCOAL), ('No data on endpoint', False, BRAND_GREY),
         ('BitLocker + remote wipe', False, BRAND_GREY), ('Equivalent', False, TEAL)],
        [('Session logging', True, DARK_CHARCOAL), ('VDI session recording', False, BRAND_GREY),
         ('CloudBrink + M365 + MyAudit', False, BRAND_GREY), ('Equivalent', False, TEAL)],
        [('User experience', True, DARK_CHARCOAL), ('Datacenter latency, no offline', False, BRAND_GREY),
         ('Native performance, offline OK', False, BRAND_GREY), ('Improved', False, TEAL)],
    ])
notes(s, """TALKING POINTS:
- Head-to-head. Walk through each row.
- Access control: device posture > IP address. Validates what the device IS, not where it is.
- Travel: managed devices win. No lockouts, no exceptions.
- Device loss: BitLocker + remote wipe compensates. Equivalent outcome.
- User experience: pure improvement, no compliance trade-off.""")


# ═════════════════════════════════════════════
# SLIDE 8: MOBILE DEVICE CONTROL
# ═════════════════════════════════════════════
s = new_slide()
chrome(s, 'Mobile Device Control', 8)
card4(s, 0, 'Intune MDM Enrollment',
      ['Company phone enrolls in Intune before M365 apps activate.',
       'Enforces: encryption, PIN/biometric, OS minimum, jailbreak detection.',
       'Unenrolled personal phones blocked.',
       'Included in M365 E3.'],
      'Reg S-P \u00a7248.30(a)(1)-(3)')
card4(s, 1, 'App Protection (MAM)',
      ['Managed container around Outlook, OneDrive, Teams.',
       'Blocked: copy/paste to personal apps, save-as, screenshots.',
       'Attachments open only in managed apps.',
       'Encrypted independently.'],
      'Reg S-P \u00a7248.30(a)(3)')
card4(s, 2, 'Remote Wipe',
      ['Selective: corporate data only. Personal untouched.',
       'Full: factory reset on termination.',
       'Auto triggers: jailbreak, days offline.',
       'Wipe history logged.'],
      'FINRA 3110(a); NIST MP-6')
card4(s, 3, 'Monthly Mobile Audit',
      ['Enrolled count vs. headcount.',
       'Compliance violations, wipe history.',
       'Evidence for FINRA 3110(c).',
       'Folds into $50/mo assessment.'],
      'FINRA 3110(c); FINRA 3120')
notes(s, """TALKING POINTS:
- All 9 users get company phones. Phone with Outlook = regulated data.
- MDM: Must enroll before M365 activates. No enrollment = no email.
- MAM: Managed container. No copy to WhatsApp, no save to personal storage.
- Remote wipe: selective (corp data only) or full. Auto on jailbreak.
- Monthly audit: added to existing assessment, no extra cost. Intune MDM/MAM in E3.""")


# ═════════════════════════════════════════════
# SLIDE 9: ENDPOINT DLP
# ═════════════════════════════════════════════
s = new_slide()
table_slide(s, 'Endpoint DLP: MyAudit + M365', 9,
    ['Data Movement Path', 'M365 E3 DLP', 'MyAudit'],
    [
        [('Download to local C: drive', True, DARK_CHARCOAL), ('Not covered', False, BRAND_GREY), ('Blocked / logged', False, BRAND_GREY)],
        [('Copy to USB drive', True, DARK_CHARCOAL), ('Not covered', False, BRAND_GREY), ('Blocked / logged', False, BRAND_GREY)],
        [('Print to local printer', True, DARK_CHARCOAL), ('Not covered', False, BRAND_GREY), ('Blocked / logged', False, BRAND_GREY)],
        [('Clipboard copy/paste', True, DARK_CHARCOAL), ('Not covered', False, BRAND_GREY), ('Blocked / logged', False, BRAND_GREY)],
        [('Screen capture', True, DARK_CHARCOAL), ('Not covered', False, BRAND_GREY), ('Blocked / logged', False, BRAND_GREY)],
        [('Send via Exchange email', True, DARK_CHARCOAL), ('Covered (mail DLP)', False, BRAND_GREY), ('N/A (M365)', False, BRAND_GREY)],
        [('SharePoint/OneDrive share', True, DARK_CHARCOAL), ('Covered (M365 DLP)', False, BRAND_GREY), ('N/A (tenant)', False, BRAND_GREY)],
        [('Upload to unsanctioned SaaS', True, DARK_CHARCOAL), ('Partial (needs E5)', False, BRAND_GREY), ('Blocked / logged', False, BRAND_GREY)],
    ])
notes(s, """TALKING POINTS:
- M365 DLP = cloud-side. MyAudit = device-side. Complementary, not duplicative.
- Without MyAudit, first five rows uncontrolled. User downloads to USB and walks out.
- Same as VDI proposal. Same agent, same policies, same blocks.
- SEC exam letters have cited firms with M365 DLP but missing endpoint DLP.""")


# ═════════════════════════════════════════════
# SLIDE 10: TENANT BACKUP
# ═════════════════════════════════════════════
s = new_slide()
chrome(s, 'Tenant Backup: Veeam for M365', 10)
card3(s, 0, 'M365 Retention Gaps',
      ['\u2022 Tenant-level events (malicious admin, takeover)',
       '\u2022 Ransomware via OneDrive/SharePoint sync',
       '\u2022 Mass deletion across tenant',
       '\u2022 Retention-policy changes',
       '\u2022 No point-in-time item/folder restore',
       'If the tenant is the attack surface, retention is not a backup.'],
      'SEC Rule 17a-4(f); FINRA 4511')
card3(s, 1, 'What Veeam Adds',
      ['\u2022 Full tenant: Exchange, SharePoint, OneDrive, Teams',
       '\u2022 Immutable off-tenant storage (WORM-capable)',
       '\u2022 Granular restore: item, folder, mailbox, site',
       '\u2022 Independent encryption keys',
       '\u2022 Evidence chain for 17a-4(f) attestation'],
      'SEC Rule 17a-4(b),(f); FINRA 4511')
card3(s, 2, 'Unchanged From VDI Proposal',
      ['Veeam backs up the M365 tenant, not the endpoint.',
       'Device model is irrelevant \u2014 same backup scope.',
       '$100/month storage = independent backup repository.',
       'This makes 17a-4(f) defensible regardless of endpoint model.'],
      'No change from original proposal')
notes(s, """TALKING POINTS:
- Unchanged from VDI proposal. Veeam backs up tenant, not endpoint.
- M365 retention is INSIDE the tenant. If tenant is compromised, retention is part of the attack surface.
- SEC 17a-4(f) requires non-rewriteable, non-erasable preservation. Microsoft retention alone doesn't meet this.
- $100/month storage is the independent backup repository.""")


# ═════════════════════════════════════════════
# SLIDE 11: DRIFT DETECTION
# ═════════════════════════════════════════════
s = new_slide()
chrome(s, 'Drift Detection: Monthly Assessment', 11)
card3(s, 0, 'Recommended: $50/Month',
      ['Monthly Technijian assessment:',
       '\u2022 M365: DLP, MFA, audit logs, privileged roles',
       '\u2022 Intune: compliance, encryption, OS patches',
       '\u2022 CloudBrink: agents vs. headcount, posture',
       '\u2022 Mobile: enrolled devices, wipe history, MAM',
       '\u2022 Physical site controls',
       'Monthly evidence packet for FINRA 3110(c).'],
      'FINRA 3110(c); FINRA 3120')
card3(s, 1, 'Alternative: $250/Quarter',
      ['Quarterly if AFFG prefers lower steady-state cost.',
       'Full engagement each cycle: scan, review, evidence.',
       'Not 3x monthly \u2014 each cycle restarts full assessment.',
       'Minimum cadence for FINRA 3110(c).'],
      'FINRA 3110(c) (minimum)')
card3(s, 2, 'New: Device Fleet Monitoring',
      ['Two new assessment dimensions:',
       '1. Intune: are all 9 laptops + 9 phones enrolled, encrypted, patched, checking in?',
       '2. CloudBrink: agents active, device checks passing, unauthorized access attempts?',
       'These replace VDI session monitoring.'],
      'Reg S-P \u00a7248.30(a); NIST CM-8')
notes(s, """TALKING POINTS:
- Same $50/mo, expanded scope: device fleet, CloudBrink, mobile audit.
- Monthly = cleanest FINRA 3110(c) evidence. Quarterly = minimum acceptable.
- New: Intune device compliance (18 devices total) and CloudBrink posture replace VDI monitoring.""")


# ═════════════════════════════════════════════
# SLIDE 12: NO DUPLICATION
# ═════════════════════════════════════════════
s = new_slide()
table_slide(s, 'No Duplication: Each Layer Solves A Distinct Problem', 12,
    ['Layer', 'Unique Control', 'If Removed...'],
    [
        [('Intune-managed laptops', True, DARK_CHARCOAL), ('Encrypted, policy-enforced endpoint', False, BRAND_GREY),
         ('Personal devices, no encryption/wipe', False, BRAND_GREY)],
        [('Conditional Access', True, DARK_CHARCOAL), ('M365 gated to compliant devices', False, BRAND_GREY),
         ('Personal devices access email via OWA/IMAP', False, BRAND_GREY)],
        [('CloudBrink ZTNA', True, DARK_CHARCOAL), ('Device-posture access, any location', False, BRAND_GREY),
         ('IP whitelist (travel blocked) or open', False, BRAND_GREY)],
        [('MyAudit (endpoint DLP)', True, DARK_CHARCOAL), ('Blocks USB, download, print, clipboard', False, BRAND_GREY),
         ('Exfiltration via USB/download', False, BRAND_GREY)],
        [('SSO / 2FA gateway', True, DARK_CHARCOAL), ('MFA on custodian portals', False, BRAND_GREY),
         ('Portal MFA unauditable', False, BRAND_GREY)],
        [('Intune MDM/MAM (phones)', True, DARK_CHARCOAL), ('Phone containerization + wipe', False, BRAND_GREY),
         ('Phone loss exposes M365 email/files', False, BRAND_GREY)],
        [('Veeam Backup', True, DARK_CHARCOAL), ('Immutable tenant backup', False, BRAND_GREY),
         ('Ransomware/deletion unrecoverable', False, BRAND_GREY)],
        [('Monthly assessment', True, DARK_CHARCOAL), ('Attested evidence of controls', False, BRAND_GREY),
         ('Drift undetected, 3110(c) gap', False, BRAND_GREY)],
    ])
notes(s, """TALKING POINTS:
- Walk through "If Removed" column. Each layer has a unique regulatory gap it closes.
- Conditional Access: without it, personal devices access email via OWA or legacy auth.
- Every layer has a cite-able FINRA/SEC requirement behind it.""")


# ═════════════════════════════════════════════
# SLIDE 13: REVISED PATH
# ═════════════════════════════════════════════
s = new_slide()
bullet_slide(s, 'What Changes From The VDI Proposal', 13,
    [
        (CORE_BLUE, 'VDI removed. Replaced by company-owned, Intune-managed laptops for all 9 users.'),
        (TEAL, 'CloudBrink ZTNA added. Replaces IP whitelist. Device-posture access from any location.'),
        (CORE_ORANGE, 'Company phones added. Intune MDM/MAM, app containerization, remote wipe for all 9 users.'),
        (CORE_BLUE, 'Conditional Access: M365 restricted to compliant devices. Legacy auth blocked.'),
        (TEAL, 'MyAudit endpoint DLP retained. Same exfiltration blocks on managed laptops.'),
        (CORE_ORANGE, 'SSO/2FA gateway retained. MFA on custodian portals and non-M365 systems.'),
        (CORE_BLUE, 'Veeam Backup retained. Full tenant backup \u2014 unchanged.'),
        (TEAL, 'Monthly assessment at $50/mo. Scope expanded: device fleet + CloudBrink + mobile.'),
        (CORE_ORANGE, 'No VDI hosting cost. One-time: Intune, CloudBrink, MyAudit, phone provisioning.'),
    ])
notes(s, """TALKING POINTS:
- REMOVED: VDI hosting.
- ADDED: CloudBrink ZTNA, company phones with Intune, Conditional Access policies.
- RETAINED: MyAudit, SSO/2FA, Veeam, monthly assessment.
- Assessment expanded at no additional cost. Regulatory posture maintained or strengthened.""")


# ═════════════════════════════════════════════
# SLIDE 14: CLOSING (Dark Charcoal)
# ═════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, SW, SH, DARK_CHARCOAL)
s.shapes.add_picture(LOGO_WHITE, Inches(3.5), Inches(0.5), Inches(4.5), Inches(1.5))
txt(s, Inches(1), Inches(2.2), Inches(11), Inches(1.0), 'The Bottom Line', Pt(40), True, WHITE, PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(3.3), Inches(10), Inches(0.8),
    'Same FINRA and SEC compliance posture \u2014 without VDI. Managed devices, CloudBrink ZTNA, and Intune close every gap the VDI model closed.',
    Pt(16), color=LIGHT_GREY, align=PP_ALIGN.CENTER)
pts = [
    'Same four-factor boundary: user + MFA + managed device + CloudBrink ZTNA',
    'Laptops AND phones under one Intune management umbrella',
    'Travel-ready: CloudBrink works from any network',
    'Every control maps to a specific FINRA or SEC citation',
]
for i, p in enumerate(pts):
    y = Inches(4.4) + i * Inches(0.5)
    rect(s, Inches(3.0), y+Inches(0.1), Inches(0.1), Inches(0.1), TEAL)
    txt(s, Inches(3.3), y, Inches(7), Inches(0.45), p, Pt(14), color=WHITE)
rect(s, 0, Inches(6.85), SW, Inches(0.08), CORE_ORANGE)
txt(s, Inches(7.0), Inches(6.4), Inches(5.5), Inches(0.3), 'technology as a solution', Pt(11), color=TEAL, align=PP_ALIGN.RIGHT)
txt(s, Inches(0.6), Inches(7.05), Inches(5), Inches(0.3),
    'Technijian  |  AFFG Managed Device Strategy  |  April 2026', Pt(10), color=BRAND_GREY)
notes(s, """TALKING POINTS:
- Four key takeaways: four-factor boundary, one Intune umbrella, travel-ready, citation-mapped.
- CloudBrink is stronger than IP whitelisting.
- Every control has a regulatory requirement behind it.
- NEXT: Iris reviews, proceed to revised proposal with pricing.""")


# ═════════════════════════════════════════════
out = 'AFFG_Managed_Device_Strategy_Technijian.pptx'
prs.save(out)
print(f'Saved {out} with {len(prs.slides)} slides')
